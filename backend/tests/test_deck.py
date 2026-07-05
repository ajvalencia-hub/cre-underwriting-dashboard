"""H12: one-page deck export — structure, pass-through numbers, and the
router's degradation path."""

import json
from io import BytesIO
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from app.database import Base, get_db
from app.main import app
from app.services import deck_service
from app.services.proforma import engine

FIXTURES = Path(__file__).parent / "fixtures"


@pytest.fixture
def analytic() -> dict:
    return json.loads((FIXTURES / "analytic_acquisition.json").read_text())


@pytest.fixture
def client():
    eng = create_engine(
        "sqlite://", connect_args={"check_same_thread": False}, poolclass=StaticPool
    )
    Base.metadata.create_all(eng)
    TestSession = sessionmaker(bind=eng)

    def _override():
        db = TestSession()
        try:
            yield db
        finally:
            db.close()

    app.dependency_overrides[get_db] = _override
    yield TestClient(app)
    app.dependency_overrides.pop(get_db)
    eng.dispose()


def _slide_text(pptx_bytes: bytes) -> str:
    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) == 1  # ONE page, that's the product
    chunks = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def test_deck_is_one_slide_with_passthrough_numbers(analytic):
    result = engine.compute(analytic)
    content = deck_service.build_deck("Analytic Acquisition", analytic, result)
    text = _slide_text(content)

    assert "Analytic Acquisition" in text
    # Levered IRR formatted exactly as the engine computed it.
    assert f"{result['outputs']['leveredIrr'] * 100:.2f}%" in text
    assert f"{result['outputs']['equityMultiple']:.2f}x" in text
    assert "KEY ASSUMPTIONS" in text
    assert "$1,000,000" in text  # purchase price
    assert "internal discussion only" in text

    # Charts embedded as pictures (cash flow + sources & uses).
    prs = Presentation(BytesIO(content))
    pictures = [s for s in prs.slides[0].shapes if s.shape_type == 13]  # PICTURE
    assert len(pictures) == 2


def test_deck_route_and_degradation(client, analytic):
    deal = client.post(
        "/api/deals", json={"name": "Deck <Deal>", "inputs": analytic}
    ).json()
    response = client.get(f"/api/deals/{deal['id']}/deck.pptx")
    assert response.status_code == 200
    assert response.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml"
    )
    assert "<" not in response.headers["content-disposition"]
    assert _slide_text(response.content)  # parses as a real pptx

    empty = client.post("/api/deals", json={"name": "Empty", "inputs": {}}).json()
    missing = client.get(f"/api/deals/{empty['id']}/deck.pptx")
    assert missing.status_code == 422
    assert "missing inputs" in missing.json()["detail"]

    assert client.get("/api/deals/nope/deck.pptx").status_code == 404
