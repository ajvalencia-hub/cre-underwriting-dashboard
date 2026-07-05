"""One-page deck export (H12) via python-pptx.

Same STRICT RULE as the memo and the HTML share: zero financial math here.
Every number is a formatted pass-through from a fresh engine compute; the
charts are the memo's own matplotlib PNGs (annual levered cash flow,
sources & uses) fed the engine's vectors.
"""

from datetime import date
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.config import FIRM_NAME, MEMO_BRAND_COLOR
from app.services import memo_charts
from app.services.memo_service import _OUTPUT_META, format_value

# Metric tiles across the top, in order (shown when present).
_TILE_IDS = [
    "leveredIrr",
    "equityMultiple",
    "cashOnCashYear1",
    "goingInCapRate",
    "minDscr",
    "npv",
]

_ASSUMPTION_ROWS = [
    ("purchasePrice", "Purchase price", "currency"),
    ("landCost", "Land cost", "currency"),
    ("hardCosts", "Hard costs", "currency"),
    ("grossPotentialRent", "Gross potential rent", "currency"),
    ("vacancyPct", "Vacancy", "percent"),
    ("rentGrowthPct", "Rent growth", "percent"),
    ("holdPeriodYears", "Hold (yrs)", "number"),
    ("exitCapRatePct", "Exit cap", "percent"),
    ("ltvOrLtc", "LTV / LTC", "percent"),
    ("interestRate", "Interest rate", "percent"),
]

_SLIDE_W = Inches(13.333)  # 16:9
_SLIDE_H = Inches(7.5)
_MARGIN = Inches(0.45)

_DISCLAIMER = (
    "Prepared for internal discussion only. Projections are estimates based on "
    "the stated assumptions; actual results will differ. Not an offer to sell "
    "or a solicitation of an offer to buy any security."
)


def _brand() -> RGBColor:
    return RGBColor.from_string(MEMO_BRAND_COLOR)


def _text(slide, left, top, width, height, text, size, *, bold=False,
          color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    if color is not None:
        font.color.rgb = color
    return box


def build_deck(deal_name: str, inputs: dict, result: dict) -> bytes:
    """One 16:9 slide: title bar, metric tiles, assumptions column, and the
    memo's cash-flow + sources & uses charts."""
    outputs = result.get("outputs", {})
    statement = result.get("statement")
    sources_and_uses = result.get("sourcesAndUses")

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    _text(slide, _MARGIN, Inches(0.25), _SLIDE_W - 2 * _MARGIN, Inches(0.5),
          deal_name, 26, bold=True, color=_brand())
    _text(slide, _MARGIN, Inches(0.72), _SLIDE_W - 2 * _MARGIN, Inches(0.3),
          f"{FIRM_NAME} · Investment summary · {date.today().isoformat()}",
          10, color=RGBColor.from_string("64748B"))

    # Metric tiles
    tiles = [(tid, _OUTPUT_META.get(tid, {})) for tid in _TILE_IDS if tid in outputs]
    if tiles:
        tile_w = Emu(int((_SLIDE_W - 2 * _MARGIN) / max(4, len(tiles))))
        top = Inches(1.15)
        for i, (tid, meta) in enumerate(tiles):
            left = _MARGIN + i * tile_w
            _text(slide, left, top, tile_w, Inches(0.25),
                  str(meta.get("label", tid)).upper(), 9,
                  color=RGBColor.from_string("94A3B8"))
            _text(slide, left, top + Inches(0.24), tile_w, Inches(0.4),
                  format_value(outputs[tid], meta.get("type", "number")),
                  20, bold=True)

    # Assumptions column (left)
    assumptions_top = Inches(2.15)
    _text(slide, _MARGIN, assumptions_top, Inches(3.6), Inches(0.25),
          "KEY ASSUMPTIONS", 10, bold=True, color=_brand())
    row_top = assumptions_top + Inches(0.35)
    for field_id, label, value_type in _ASSUMPTION_ROWS:
        value = inputs.get(field_id)
        if value in (None, "", 0):
            continue
        _text(slide, _MARGIN, row_top, Inches(2.0), Inches(0.22), label, 10,
              color=RGBColor.from_string("64748B"))
        _text(slide, _MARGIN + Inches(2.0), row_top, Inches(1.6), Inches(0.22),
              format_value(value, value_type), 10, align=PP_ALIGN.RIGHT)
        row_top += Inches(0.28)
        if row_top > Inches(6.5):
            break

    # Charts (right two-thirds)
    chart_left = Inches(4.5)
    chart_w = _SLIDE_W - chart_left - _MARGIN
    cashflow_png = memo_charts.annual_cashflow_bars(statement)
    if cashflow_png:
        slide.shapes.add_picture(BytesIO(cashflow_png), chart_left, Inches(2.15),
                                 width=chart_w)
    sources_png = memo_charts.sources_uses_bars(sources_and_uses)
    if sources_png:
        slide.shapes.add_picture(BytesIO(sources_png), chart_left, Inches(4.55),
                                 width=chart_w)

    # Footer
    _text(slide, _MARGIN, _SLIDE_H - Inches(0.45), _SLIDE_W - 2 * _MARGIN,
          Inches(0.35), _DISCLAIMER, 7, color=RGBColor.from_string("94A3B8"))

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
